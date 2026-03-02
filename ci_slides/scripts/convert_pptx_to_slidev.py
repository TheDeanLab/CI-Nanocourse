#!/usr/bin/env python3
"""Convert lecture PPTX files to cleaned Slidev markdown."""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation


COMMAND_PREFIXES = (
    "git ",
    "conda ",
    "pip ",
    "python ",
    "pytest ",
    "npm ",
    "pnpm ",
    "module ",
)

EXACT_REPLACEMENTS = {
    "And it’s Role in Facilitating Team-Based Development.": "And its role in facilitating team-based development.",
}


@dataclass
class SlideData:
    title: str = ""
    subtitle: str = ""
    bullets: list[tuple[int, str]] = field(default_factory=list)


def normalize(text: str) -> str:
    value = text.strip().lower()
    value = value.replace("’", "'").replace("“", '"').replace("”", '"')
    value = re.sub(r"\s+", " ", value)
    value = re.sub(r"[`*_]", "", value)
    return value


def is_intermission(slide: SlideData) -> bool:
    if normalize(slide.title) == "intermission":
        return True
    if normalize(slide.subtitle) == "intermission":
        return True
    return len(slide.bullets) == 1 and normalize(slide.bullets[0][1]) == "intermission"


def is_activity(slide: SlideData) -> bool:
    joined = " ".join([slide.title, slide.subtitle, *(t for _, t in slide.bullets)])
    return "activity #" in normalize(joined)


def looks_command(text: str) -> bool:
    lowered = text.strip().lower()
    if any(lowered.startswith(prefix) for prefix in COMMAND_PREFIXES):
        return True
    return bool(re.match(r"^\d+[a-z]?\.\s*(git|conda|pip|python|npm|pnpm)\b", lowered))


def format_text(text: str) -> str:
    value = text.strip()
    value = EXACT_REPLACEMENTS.get(value, value)
    value = re.sub(r"\bconfilcts\b", "conflicts", value, flags=re.IGNORECASE)
    value = re.sub(r"\bobsolte\b", "obsolete", value, flags=re.IGNORECASE)
    value = re.sub(r"\bseaborns\b", "seaborn", value, flags=re.IGNORECASE)
    value = re.sub(r"\bGIT\b", "Git", value)
    if "{{" in value or "}}" in value:
        return f"`{value}`"
    if re.fullmatch(r"https?://\S+", value):
        return f"[{value}]({value})"
    return value


def extract_slide(slide) -> SlideData:
    title = ""
    if slide.shapes.title is not None and slide.shapes.title.text:
        title = slide.shapes.title.text.strip()

    raw_lines: list[tuple[int, str]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if not text:
                continue
            level = int(paragraph.level or 0)
            raw_lines.append((level, text))

    filtered: list[tuple[int, str]] = []
    seen: set[tuple[int, str]] = set()
    title_norm = normalize(title)
    for level, text in raw_lines:
        text_norm = normalize(text)
        if not text_norm:
            continue
        if text_norm == title_norm:
            continue
        key = (level, text_norm)
        if key in seen:
            continue
        seen.add(key)
        filtered.append((level, text))

    subtitle = ""
    bullets = filtered
    if filtered:
        first_level, first_text = filtered[0]
        first_norm = normalize(first_text)
        if (
            first_level == 0
            and len(first_text) <= 100
            and not looks_command(first_text)
            and not first_norm.endswith(":")
            and not re.fullmatch(r"https?://\S+", first_text.strip())
        ):
            subtitle = first_text
            bullets = filtered[1:]

    return SlideData(title=title, subtitle=subtitle, bullets=bullets)


def should_merge(previous: SlideData, current: SlideData) -> bool:
    if is_activity(previous) or is_activity(current):
        return False
    if is_intermission(previous) or is_intermission(current):
        return False
    if normalize(previous.title) != normalize(current.title):
        return False
    if normalize(previous.subtitle) != normalize(current.subtitle):
        return False
    if len(previous.bullets) + len(current.bullets) > 8:
        return False
    return bool(previous.title or previous.subtitle)


def merge_slide(previous: SlideData, current: SlideData) -> SlideData:
    existing = {normalize(text) for _, text in previous.bullets}
    for level, text in current.bullets:
        key = normalize(text)
        if key in existing:
            continue
        previous.bullets.append((level, text))
        existing.add(key)
    return previous


def collect_slides(pptx_path: Path) -> list[SlideData]:
    presentation = Presentation(pptx_path)
    slides = [extract_slide(slide) for slide in presentation.slides]
    merged: list[SlideData] = []
    for slide in slides:
        if not slide.title and not slide.subtitle and not slide.bullets:
            continue
        if merged and should_merge(merged[-1], slide):
            merged[-1] = merge_slide(merged[-1], slide)
        else:
            merged.append(slide)
    return merged


def render_slide(slide: SlideData) -> str:
    if is_intermission(slide):
        return (
            "---\n"
            "layout: center\n"
            "class: text-center\n"
            "---\n"
            "# Intermission\n"
        )

    content: list[str] = ["---"]

    heading = slide.title.strip()
    subtitle = slide.subtitle.strip()

    if not heading and subtitle and not slide.bullets:
        content.extend(
            [
                "layout: cover",
                "class: text-center",
                "---",
                f"# {format_text(subtitle)}",
            ]
        )
        return "\n".join(content) + "\n"

    if heading:
        content.append(f"# {format_text(heading)}")
    elif subtitle:
        content.append(f"# {format_text(subtitle)}")
        subtitle = ""

    if subtitle:
        content.append(f"## {format_text(subtitle)}")

    if slide.bullets:
        content.append("")
        min_level = min(level for level, _ in slide.bullets)
        for level, text in slide.bullets:
            normalized_level = max(level - min_level, 0)
            indent = "  " * normalized_level
            content.append(f"{indent}- {format_text(text)}")

    return "\n".join(content) + "\n"


def render_lecture(lecture_name: str, slides: list[SlideData]) -> str:
    output: list[str] = [f"<!-- {lecture_name} -->\n"]
    for slide in slides:
        output.append(render_slide(slide))
    return "\n".join(part.rstrip() for part in output).rstrip() + "\n"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert lecture PPTX files to Slidev markdown.")
    parser.add_argument(
        "--source-root",
        type=Path,
        default=Path("source_slides"),
        help="Directory containing lecture-XX folders with lecture-XX.pptx files.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=Path("lectures"),
        help="Directory where generated lecture markdown files will be written.",
    )
    parser.add_argument(
        "--lectures",
        nargs="*",
        default=[f"lecture-{i:02d}" for i in range(1, 10)],
        help="Lecture folder names to convert.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output_root = args.output_root
    output_root.mkdir(parents=True, exist_ok=True)

    for lecture in args.lectures:
        lecture_dir = args.source_root / lecture
        pptx_path = lecture_dir / f"{lecture}.pptx"
        if not pptx_path.exists():
            raise FileNotFoundError(f"Missing PPTX file: {pptx_path}")

        slides = collect_slides(pptx_path)
        rendered = render_lecture(lecture, slides)
        out_path = output_root / f"{lecture}.md"
        out_path.write_text(rendered, encoding="utf-8")
        print(f"Wrote {out_path} ({len(slides)} slides)")


if __name__ == "__main__":
    main()
