#!/usr/bin/env python3
"""Apply requested media updates for lectures 4-9.

This script:
1. Adds image-right layouts to requested slides in lectures 4 and 5.
2. Adds image-right layouts for all image-bearing content in lectures 7, 8, and 9.
3. Inserts code blocks for requested lecture-06 pages.
"""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from math import ceil
from pathlib import Path
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path("/Users/Dean/Documents/GitHub/teaching/CI-Nanocourse/ci_slides")
SOURCE_ROOT = ROOT / "source_slides"
LECTURE_ROOT = ROOT / "lectures"
PUBLIC_IMG_ROOT = ROOT / "public" / "images"

PHOTO_START = "<!-- photos:start -->"
PHOTO_END = "<!-- photos:end -->"
CODE_START = "<!-- code:start -->"
CODE_END = "<!-- code:end -->"


@dataclass(frozen=True)
class ImageTask:
    lecture_file: str
    source_lecture: str
    source_slides: tuple[int, ...]
    h1: Optional[str] = None
    h2: Optional[str] = None
    contains: Optional[str] = None
    occurrence: int = 1
    target_index: Optional[int] = None  # 1-based


@dataclass(frozen=True)
class CodeTask:
    target_index: int  # 1-based
    language: str
    code: str


IMAGE_TASKS: list[ImageTask] = [
    ImageTask("lecture-04.md", "lecture-04", (1,), h1="Repository Organizational Strategies", h2="Fighting Entropy"),
    ImageTask("lecture-04.md", "lecture-04", (2,), h1="Repository Organizational Strategies", h2="Must Haves", occurrence=1),
    ImageTask("lecture-04.md", "lecture-04", (3,), h1="Repository Organizational Strategies", h2="Must Haves", occurrence=2),
    ImageTask("lecture-04.md", "lecture-04", (4,), h1="Repository Organizational Strategies", h2="Bonus Tools"),
    ImageTask("lecture-04.md", "lecture-04", (7,), h1="Model-View-Controller Architecture", h2="A Guiding Framework"),
    ImageTask("lecture-05.md", "lecture-05", (3,), contains="atlassian.com/continuous-delivery"),
    ImageTask("lecture-05.md", "lecture-05", (4,), contains="deploybot.com/blog/continuous-development"),
    ImageTask("lecture-05.md", "lecture-05", (9,), h1="How do we know this function works?"),
    ImageTask("lecture-05.md", "lecture-05", (10,), h1="We check it with a unit test"),
    ImageTask(
        "lecture-05.md",
        "lecture-05",
        (16,),
        h1="The goal of CI/CD is to ensure developments do not stray far from the main branch",
    ),
    ImageTask("lecture-05.md", "lecture-05", (22,), h1="GitHub Actions Dashboard"),
    ImageTask("lecture-05.md", "lecture-05", (23,), h1="GitHub Actions Workflow Example"),
    ImageTask(
        "lecture-07.md",
        "lecture-07",
        (4,),
        h1="Types of Tests",
        h2="Unit Test: Isolated method, function or component.",
    ),
    ImageTask("lecture-07.md", "lecture-07", (7,), h1="Running pytest", h2="Just type >> pytest in prompt"),
    ImageTask("lecture-07.md", "lecture-07", (8,), h1="Running pytest", h2="Assert is False, the test fails!"),
    ImageTask("lecture-07.md", "lecture-07", (9,), h1="Any “test_” function will be recognized by pytest automatically"),
    ImageTask("lecture-07.md", "lecture-07", (10,), h1="Parameterization in tests"),
    ImageTask("lecture-07.md", "lecture-07", (12,), h1="Fixtures in pytest", occurrence=1),
    ImageTask("lecture-07.md", "lecture-07", (13,), h1="Fixtures in pytest", occurrence=2),
    ImageTask("lecture-07.md", "lecture-07", (14,), h1="Key Binding Test in PyCalc Exercise", h2="conftest.py"),
    ImageTask("lecture-07.md", "lecture-07", (15,), h1="What is yield doing?", h2="SETUP: prepare resources"),
    ImageTask(
        "lecture-08.md",
        "lecture-08",
        (1,),
        h1="Setting up a CI/CD Pipeline in GitHub Actions",
        h2="Introduction to Python Software Development on GitHub 2023",
    ),
    ImageTask("lecture-08.md", "lecture-08", (2,), h1="GitHub Actions Dashboard"),
    ImageTask("lecture-08.md", "lecture-08", (3, 4), h1="GitHub Actions Workflow Example"),
    ImageTask("lecture-08.md", "lecture-08", (6,), h1="Where does the workflow file go?"),
    ImageTask("lecture-08.md", "lecture-08", (7,), h1="github-actions-demo.yml"),
    ImageTask("lecture-08.md", "lecture-08", (8, 9), h1="What are these two lines doing??"),
    ImageTask("lecture-08.md", "lecture-08", (11,), h1="What does               do?"),
    ImageTask("lecture-08.md", "lecture-08", (12,), h1="Defining the jobs", h2="What is jobs?"),
    ImageTask("lecture-08.md", "lecture-08", (13, 14, 15, 16, 17, 18, 19), h1="Defining the jobs", occurrence=2),
    ImageTask("lecture-08.md", "lecture-08", (20, 21), h1="The full Explore-GitHub-Actions job"),
    ImageTask("lecture-08.md", "lecture-08", (23,), h1="GitHub Marketplace"),
    ImageTask("lecture-08.md", "lecture-08", (24,), h1="Search for Codecov"),
    ImageTask("lecture-08.md", "lecture-08", (25,), h1="Click “Use latest version”"),
    ImageTask("lecture-08.md", "lecture-08", (26,), h1="Copy the action as a step in your job"),
    ImageTask("lecture-08.md", "lecture-08", (27,), h1="There’s an action for everything"),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (1,),
        h1="Public-facing documentation",
        h2="Introduction to Python Software Development on GitHub 2023",
    ),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (4,),
        h1="Software licenses are an important part of documentation, and should be chosen based on what you want to do with your software",
    ),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (5,),
        h1="GPLv3 is the premier copyleft license, but if you want to license your code to a company, you are better off with an MIT license",
    ),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (8,),
        h1="Documentation languages are simple languages, close to standard word processing",
    ),
    ImageTask("lecture-09.md", "lecture-09", (9,), h1="GitHub uses Markdown"),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (11,),
        h1="Python docstrings can be written in a variety of documentation languages, including reStructuredText",
        h2="reStructuredText",
    ),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (12,),
        h1="Python docstrings can be written in a variety of documentation languages, including reStructuredText",
        h2="NumPy style",
    ),
    ImageTask("lecture-09.md", "lecture-09", (14,), h1="is the primary documentation framework for Python"),
    ImageTask("lecture-09.md", "lecture-09", (15, 16), h1="An example of a docstring compiled to HTML by Sphinx"),
    ImageTask("lecture-09.md", "lecture-09", (19, 20), h1="Let’s generate a source folder"),
    ImageTask("lecture-09.md", "lecture-09", (21,), h1="Directory structure of documentation folder"),
    ImageTask("lecture-09.md", "lecture-09", (22,), h1="Now try make html in the docs folder"),
    ImageTask("lecture-09.md", "lecture-09", (23,), h1="The toctree defines what we see"),
    ImageTask("lecture-09.md", "lecture-09", (24, 25), h1="Autosummary helps us with the API documentation"),
    ImageTask("lecture-09.md", "lecture-09", (26,), h1="Templating"),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (27,),
        h1="Templates enable us to produce a more comprehensive autosummary",
        occurrence=1,
    ),
    ImageTask(
        "lecture-09.md",
        "lecture-09",
        (28,),
        h1="Templates enable us to produce a more comprehensive autosummary",
        occurrence=2,
    ),
    ImageTask("lecture-09.md", "lecture-09", (29,), h1="Exercises"),
    ImageTask("lecture-09.md", "lecture-09", (32,), h1="GitHub Pages"),
]


CODE_TASKS_LECTURE06: list[CodeTask] = [
    CodeTask(
        8,
        "python",
        """def process_large_dataset(data):
    \"\"\"Processes a dataset and returns useful statistics.\"\"\"
    # implementation here
    return {}


def process_large_dataset_with_very_long_name(
    data,
    additional_parameters,
    more_params,
):
    \"\"\"Prefer wrapped signatures over single long lines.\"\"\"
    return {}""",
    ),
    CodeTask(
        9,
        "python",
        """import os
import sys

x = 5
y = 10
print(x + y)""",
    ),
    CodeTask(
        11,
        "python",
        """def divide(a: float, b: float) -> float:
    \"\"\"
    Divides two numbers.

    Parameters
    ----------
    a : float
        Numerator.
    b : float
        Denominator.

    Returns
    -------
    float
        Result of division.

    Raises
    ------
    ZeroDivisionError
        If `b` is 0.
    \"\"\"
    if b == 0:
        raise ZeroDivisionError(\"Cannot divide by zero.\")
    return a / b""",
    ),
    CodeTask(
        13,
        "python",
        """from functools import wraps


def log_runtime(fn):
    @wraps(fn)
    def wrapper(*args, **kwargs):
        print(f\"Running {fn.__name__}...\")
        return fn(*args, **kwargs)

    return wrapper


@log_runtime
def load_dataset(path):
    return open(path, encoding=\"utf-8\").read()""",
    ),
    CodeTask(
        14,
        "python",
        """from pathlib import Path

config_path = Path(\"settings.toml\")

with config_path.open(\"r\", encoding=\"utf-8\") as f:
    config_text = f.read()""",
    ),
    CodeTask(
        15,
        "python",
        """def iter_valid_rows(rows):
    for row in rows:
        if row.get(\"is_valid\"):
            yield row


for row in iter_valid_rows(records):
    print(row[\"id\"])""",
    ),
    CodeTask(
        18,
        "bash",
        """black src tests
black --check .
black --diff .""",
    ),
    CodeTask(
        19,
        "bash",
        """ruff check .
ruff check . --fix""",
    ),
    CodeTask(
        21,
        "toml",
        """[tool.ruff]
line-length = 88
target-version = \"py39\"

[tool.ruff.lint]
select = [\"E\", \"F\", \"I\", \"UP\", \"D\"]
ignore = [\"D203\", \"D213\"]""",
    ),
    CodeTask(
        23,
        "yaml",
        """repos:
  - repo: https://github.com/astral-sh/ruff-pre-commit
    rev: v0.13.0
    hooks:
      - id: ruff
      - id: ruff-format
  - repo: https://github.com/psf/black
    rev: 25.1.0
    hooks:
      - id: black""",
    ),
    CodeTask(
        28,
        "bash",
        """pip install pre-commit
pre-commit sample-config > .pre-commit-config.yaml""",
    ),
    CodeTask(
        29,
        "yaml",
        """repos:
  - repo: https://github.com/astral-sh/ruff-pre-commit
    rev: v0.13.0
    hooks:
      - id: ruff
      - id: ruff-format""",
    ),
    CodeTask(
        30,
        "bash",
        """pre-commit install
pre-commit run --all-files
pre-commit run ruff --all-files""",
    ),
]


def normalize(text: str) -> str:
    return " ".join(text.replace("’", "'").replace("“", '"').replace("”", '"').lower().split())


def split_slides(markdown: str) -> list[str]:
    return markdown.split("\n---\n")


def parse_headings(block: str) -> tuple[str, str]:
    h1 = ""
    h2 = ""
    for line in block.splitlines():
        stripped = line.strip()
        if stripped.startswith("# ") and not h1:
            h1 = stripped[2:].strip()
        elif stripped.startswith("## ") and not h2:
            h2 = stripped[3:].strip()
    return h1, h2


def strip_photo_block(block: str) -> str:
    lines = block.splitlines()
    out = []
    skip = False
    for line in lines:
        stripped = line.strip()
        if stripped == PHOTO_START:
            skip = True
            continue
        if stripped == PHOTO_END:
            skip = False
            continue
        if not skip:
            out.append(line)
    return "\n".join(out).rstrip()


def strip_code_block(block: str) -> str:
    lines = block.splitlines()
    out = []
    skip = False
    for line in lines:
        stripped = line.strip()
        if stripped == CODE_START:
            skip = True
            continue
        if stripped == CODE_END:
            skip = False
            continue
        if not skip:
            out.append(line)
    return "\n".join(out).rstrip()


def strip_leading_frontmatter(block: str) -> str:
    if not block.startswith("---\n"):
        return block
    end = block.find("\n---\n", 4)
    if end == -1:
        return block
    return block[end + 5 :].lstrip("\n")


def apply_image_right_frontmatter(block: str, image_path: str) -> str:
    body = strip_photo_block(block)
    body = strip_leading_frontmatter(body)
    frontmatter = (
        "---\n"
        "layout: image-right\n"
        f"image: {image_path}\n"
        "backgroundSize: contain\n"
        "---\n"
    )
    return (frontmatter + body.lstrip("\n")).rstrip()


def inject_codeblock(block: str, language: str, code: str) -> str:
    body = strip_code_block(block)
    lines = body.splitlines()
    insert_at = 0
    for i, line in enumerate(lines):
        stripped = line.strip()
        if stripped.startswith("# "):
            insert_at = i + 1
            continue
        if stripped.startswith("## "):
            insert_at = i + 1
            continue
        if stripped == "":
            if i <= insert_at:
                insert_at = i + 1
                continue
        break

    snippet = [
        CODE_START,
        f"```{language}",
        code.rstrip(),
        "```",
        CODE_END,
    ]

    out = []
    out.extend(lines[:insert_at])
    if out and out[-1].strip():
        out.append("")
    out.extend(snippet)
    if lines[insert_at:]:
        out.append("")
        out.extend(lines[insert_at:])
    return "\n".join(out).rstrip()


def find_target_index(slides: list[str], task: ImageTask) -> int:
    if task.target_index is not None:
        idx = task.target_index - 1
        if idx < 0 or idx >= len(slides):
            raise ValueError(f"target_index out of range for {task.lecture_file}: {task.target_index}")
        return idx

    matches = []
    n_h1 = normalize(task.h1 or "")
    n_h2 = normalize(task.h2 or "")
    n_contains = normalize(task.contains or "")

    for i, slide in enumerate(slides):
        h1, h2 = parse_headings(slide)
        if n_h1 and normalize(h1) != n_h1:
            continue
        if n_h2 and normalize(h2) != n_h2:
            continue
        if n_contains and n_contains not in normalize(slide):
            continue
        matches.append(i)

    if len(matches) < task.occurrence:
        raise ValueError(
            f"Unable to match slide in {task.lecture_file}: h1={task.h1!r}, h2={task.h2!r}, contains={task.contains!r}, occurrence={task.occurrence}"
        )
    return matches[task.occurrence - 1]


def extract_slide_images(source_lecture: str, slide_no: int) -> list[Image.Image]:
    pptx_path = SOURCE_ROOT / source_lecture / f"{source_lecture}.pptx"
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_no - 1]

    images = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        blob = None
        ext = "png"
        try:
            img = shape.image
            blob = img.blob
            ext = img.ext.lower()
        except Exception:
            try:
                rid = shape._pic.blipFill.blip.rEmbed  # type: ignore[attr-defined]
                part = slide.part.related_part(rid)
                if hasattr(part, "blob"):
                    blob = part.blob
                    suffix = Path(str(part.partname)).suffix
                    if suffix:
                        ext = suffix.lstrip(".").lower()
            except Exception:
                blob = None
        if blob is None:
            continue

        with Image.open(BytesIO(blob)) as im:
            converted = im.convert("RGB")
            if ext in {"tif", "tiff"}:
                # Force standard format for browser compatibility.
                buf = BytesIO()
                converted.save(buf, format="PNG")
                with Image.open(BytesIO(buf.getvalue())) as png_im:
                    converted = png_im.convert("RGB")
            images.append(converted)

    return images


def build_collage(images: list[Image.Image]) -> Image.Image:
    if len(images) == 1:
        return images[0]

    normalized_imgs = []
    for img in images:
        dup = img.copy()
        dup.thumbnail((1600, 1200))
        normalized_imgs.append(dup)

    count = len(normalized_imgs)
    cols = 1 if count == 1 else 2 if count <= 4 else 3
    rows = ceil(count / cols)
    pad = 20

    tile_w = max(im.width for im in normalized_imgs)
    tile_h = max(im.height for im in normalized_imgs)
    canvas_w = cols * tile_w + (cols + 1) * pad
    canvas_h = rows * tile_h + (rows + 1) * pad

    canvas = Image.new("RGB", (canvas_w, canvas_h), "white")
    for i, im in enumerate(normalized_imgs):
        r = i // cols
        c = i % cols
        x0 = pad + c * (tile_w + pad)
        y0 = pad + r * (tile_h + pad)
        x = x0 + (tile_w - im.width) // 2
        y = y0 + (tile_h - im.height) // 2
        canvas.paste(im, (x, y))
    return canvas


def save_task_image(task: ImageTask, target_slide_idx: int) -> str:
    imgs: list[Image.Image] = []
    for source_slide in task.source_slides:
        imgs.extend(extract_slide_images(task.source_lecture, source_slide))

    if not imgs:
        return ""

    collage = build_collage(imgs)
    out_dir = PUBLIC_IMG_ROOT / task.source_lecture
    out_dir.mkdir(parents=True, exist_ok=True)
    source_tag = "-".join(f"{s:02d}" for s in task.source_slides)
    out_name = f"t{target_slide_idx+1:02d}_from_s{source_tag}.png"
    out_path = out_dir / out_name
    collage.save(out_path, format="PNG")
    return f"/images/{task.source_lecture}/{out_name}"


def apply_image_tasks(tasks: list[ImageTask]) -> None:
    by_file: dict[str, list[ImageTask]] = {}
    for task in tasks:
        by_file.setdefault(task.lecture_file, []).append(task)

    for lecture_file, file_tasks in by_file.items():
        lecture_path = LECTURE_ROOT / lecture_file
        slides = split_slides(lecture_path.read_text(encoding="utf-8"))

        for task in file_tasks:
            idx = find_target_index(slides, task)
            image_path = save_task_image(task, idx)
            if not image_path:
                print(
                    f"warning: no extractable images for {task.source_lecture} source slides {task.source_slides}; skipped {lecture_file} slide {idx+1}"
                )
                continue
            slides[idx] = apply_image_right_frontmatter(slides[idx], image_path)
            print(
                f"updated {lecture_file} slide={idx+1} with image-right from {task.source_lecture} slides {task.source_slides}"
            )

        lecture_path.write_text("\n---\n".join(slides).rstrip() + "\n", encoding="utf-8")


def apply_code_tasks_lecture06() -> None:
    lecture_path = LECTURE_ROOT / "lecture-06.md"
    slides = split_slides(lecture_path.read_text(encoding="utf-8"))

    for task in CODE_TASKS_LECTURE06:
        idx = task.target_index - 1
        if idx < 0 or idx >= len(slides):
            raise ValueError(f"lecture-06 slide index out of range: {task.target_index}")
        slides[idx] = inject_codeblock(slides[idx], task.language, task.code)
        print(f"inserted codeblock in lecture-06.md slide={task.target_index}")

    lecture_path.write_text("\n---\n".join(slides).rstrip() + "\n", encoding="utf-8")


def refresh_master_slides() -> None:
    lectures = [
        ("01", "Introduction to Collaborative Software Development"),
        ("02", "Environment Management"),
        ("03", "Git Essentials"),
        ("04", "Organization and MVC"),
        ("05", "CI/CD Fundamentals"),
        ("06", "Code Quality and Linting"),
        ("07", "Unit Testing and TDD"),
        ("08", "GitHub Actions"),
        ("09", "Public-Facing Documentation"),
    ]

    parts = []
    parts.append(
        "---\n"
        "theme: default\n"
        "title: CI Nanocourse\n"
        "info: Lectures 1-9 converted from Keynote/PPTX source material.\n"
        "mdc: true\n"
        "---\n\n"
    )
    parts.append(
        "---\n"
        "layout: cover\n"
        "class: text-center\n"
        "---\n"
        "# CI Nanocourse\n"
        "## Lectures 1-9 (Slidev Conversion Draft)\n\n"
    )

    for number, subtitle in lectures:
        parts.append(
            "---\n"
            "layout: center\n"
            "class: text-center\n"
            "---\n"
            f"# Lecture {int(number)}\n"
            f"## {subtitle}\n\n"
        )
        lecture_text = (LECTURE_ROOT / f"lecture-{number}.md").read_text(encoding="utf-8").rstrip() + "\n\n"
        parts.append(lecture_text)

    (ROOT / "slides.md").write_text("".join(parts).rstrip() + "\n", encoding="utf-8")
    print("refreshed slides.md from lecture files")


def main() -> None:
    apply_image_tasks(IMAGE_TASKS)
    apply_code_tasks_lecture06()
    refresh_master_slides()


if __name__ == "__main__":
    main()
