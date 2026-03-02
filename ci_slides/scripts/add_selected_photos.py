#!/usr/bin/env python3
"""Extract selected PPTX photos and inject them into matching Slidev slides."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path("/Users/Dean/Documents/GitHub/teaching/CI-Nanocourse/ci_slides")
SOURCE_ROOT = ROOT / "source_slides"
LECTURE_ROOT = ROOT / "lectures"
PUBLIC_IMG_ROOT = ROOT / "public" / "images"

PHOTO_START = "<!-- photos:start -->"
PHOTO_END = "<!-- photos:end -->"


@dataclass(frozen=True)
class Target:
    lecture_file: str
    source_lecture: str
    source_slide: int
    h1: str
    h2: Optional[str] = None
    contains: Optional[str] = None
    occurrence: int = 1


TARGETS = [
    Target("lecture-01.md", "lecture-01", 2, "Introduction", "About Us"),
    Target("lecture-01.md", "lecture-01", 12, "Basic Git Commands and Workflows", contains="Git Up, Git Out, and Git Something"),
    Target("lecture-01.md", "lecture-01", 13, "Basic Git Commands and Workflows", "Cloning a Repo"),
    Target("lecture-01.md", "lecture-01", 14, "Basic Git Commands and Workflows", "Isolating Development"),
    Target("lecture-01.md", "lecture-01", 16, "Basic Git Commands and Workflows", contains="Creating a New Branch"),
    Target("lecture-01.md", "lecture-01", 17, "Basic Git Commands and Workflows", "Switching to a Branch"),
    Target("lecture-01.md", "lecture-01", 18, "Basic Git Commands and Workflows", "Merging a Branch"),
    Target("lecture-01.md", "lecture-01", 20, "Basic Git Commands and Workflows", "Cheat to Compete with GitHub Desktop"),
    Target("lecture-02.md", "lecture-02", 2, "Environment Management with Anaconda", "Creating Reproducible Development Environments"),
    Target("lecture-02.md", "lecture-02", 7, "Environment Management with Anaconda", "Activating Your New Environment"),
    Target("lecture-02.md", "lecture-02", 10, "Environment Management with Anaconda", "Creating a New Environment from a Text File"),
    Target("lecture-02.md", "lecture-02", 13, "Environment Management with Anaconda", "Creating a New Environment from a YAML File"),
    Target("lecture-02.md", "lecture-02", 14, "Advanced Environment Management", "Creating a New Environment from a pyproject.toml File"),
    Target("lecture-03.md", "lecture-03", 1, "Git Essentials", "When git push comes to shove"),
    Target("lecture-03.md", "lecture-03", 2, "Git Essentials", "Release Branching"),
    Target("lecture-03.md", "lecture-03", 3, "Git Essentials", "Feature Branching"),
    Target("lecture-03.md", "lecture-03", 4, "Git Essentials", "Trunk Based Development"),
    Target("lecture-03.md", "lecture-03", 5, "Git Essentials", "Fear of Commitment."),
    Target("lecture-03.md", "lecture-03", 6, "Git Essentials", "Pull/Merge Requests"),
    Target("lecture-03.md", "lecture-03", 7, "Git Essentials", "Rebase vs Merge"),
    Target("lecture-03.md", "lecture-03", 8, "Git Essentials", "Advanced Pull Requests"),
    Target("lecture-03.md", "lecture-03", 9, "Git Essentials", "Resolving Conflicts"),
    Target("lecture-03.md", "lecture-03", 10, "Git Essentials", "Placing “Blame”"),
    Target("lecture-03.md", "lecture-03", 12, "Git Essentials", "Issue Tracking"),
    Target("lecture-03.md", "lecture-03", 13, "Git Essentials", "Issue Creation", occurrence=1),
    Target("lecture-03.md", "lecture-03", 14, "Git Essentials", "Issue Creation", occurrence=2),
    Target("lecture-03.md", "lecture-03", 16, "Git Essentials", "Issue Formatting"),
    Target("lecture-03.md", "lecture-03", 17, "Git Essentials", "Branch Linking"),
]


def extract_images(source_lecture: str, slide_number: int) -> list[str]:
    pptx_path = SOURCE_ROOT / source_lecture / f"{source_lecture}.pptx"
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_number - 1]

    images = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        blob = None
        ext = "png"
        try:
            img = shape.image
            ext = img.ext.lower()
            blob = img.blob
        except Exception:
            # Fallback for picture-like parts that python-pptx does not expose as ImagePart.
            try:
                rid = shape._pic.blipFill.blip.rEmbed  # type: ignore[attr-defined]
                part = slide.part.related_part(rid)
                if hasattr(part, "blob"):
                    blob = part.blob
                    suffix = Path(str(part.partname)).suffix
                    if suffix:
                        ext = suffix.lstrip(".").lower()
            except Exception:
                blob = None
        if blob is None:
            continue
        area = shape.width * shape.height
        images.append((area, ext, blob))

    if not images:
        return []

    images.sort(key=lambda x: x[0], reverse=True)
    out_dir = PUBLIC_IMG_ROOT / source_lecture
    out_dir.mkdir(parents=True, exist_ok=True)
    for stale in out_dir.glob(f"s{slide_number:02d}_img*.*"):
        stale.unlink()

    paths = []
    for i, (_, ext, blob) in enumerate(images, start=1):
        if ext in {"tif", "tiff"}:
            with Image.open(BytesIO(blob)) as img:
                converted = BytesIO()
                img.save(converted, format="PNG")
                blob = converted.getvalue()
            ext = "png"

        filename = f"s{slide_number:02d}_img{i:02d}.{ext}"
        out_path = out_dir / filename
        out_path.write_bytes(blob)
        paths.append(f"/images/{source_lecture}/{filename}")
    return paths


def split_slides(markdown: str) -> list[str]:
    return markdown.split("\n---\n")


def slide_headings(block: str) -> tuple[str, str]:
    h1 = ""
    h2 = ""
    for line in block.splitlines():
        stripped = line.strip()
        if stripped.startswith("# ") and not h1:
            h1 = stripped[2:].strip()
        elif stripped.startswith("## ") and not h2:
            h2 = stripped[3:].strip()
    return h1, h2


def strip_existing_photo_block(block: str) -> str:
    lines = block.splitlines()
    out = []
    skip = False
    for line in lines:
        if line.strip() == PHOTO_START:
            skip = True
            continue
        if line.strip() == PHOTO_END:
            skip = False
            continue
        if not skip:
            out.append(line)
    return "\n".join(out).rstrip()


def render_photo_block(image_paths: list[str], alt_prefix: str) -> str:
    lines = [PHOTO_START, '<div class="mt-4 flex flex-wrap items-center justify-center gap-3">']
    for i, path in enumerate(image_paths, start=1):
        lines.append(
            f'  <img src="{path}" alt="{alt_prefix} image {i}" class="max-h-44 w-auto object-contain rounded shadow" />'
        )
    lines.append("</div>")
    lines.append(PHOTO_END)
    return "\n".join(lines)


def find_slide_index(slides: list[str], target: Target) -> int:
    matches = []
    for idx, slide in enumerate(slides):
        h1, h2 = slide_headings(slide)
        if h1 != target.h1:
            continue
        if target.h2 is not None and h2 != target.h2:
            continue
        if target.contains and target.contains not in slide:
            continue
        matches.append(idx)

    if len(matches) < target.occurrence:
        raise ValueError(
            f"No matching slide for {target.lecture_file}: h1='{target.h1}', h2='{target.h2}', contains='{target.contains}', occurrence={target.occurrence}"
        )
    return matches[target.occurrence - 1]


def inject_photos(slide: str, image_paths: list[str], alt_prefix: str) -> str:
    cleaned = strip_existing_photo_block(slide)
    lines = cleaned.splitlines()
    insert_at = 0
    for i, line in enumerate(lines):
        s = line.strip()
        if s.startswith("# "):
            insert_at = i + 1
            continue
        if s.startswith("## "):
            insert_at = i + 1
            continue
        if s == "":
            if i <= insert_at:
                insert_at = i + 1
                continue
        break

    block = render_photo_block(image_paths, alt_prefix)
    prefix = lines[:insert_at]
    suffix = lines[insert_at:]
    out = []
    out.extend(prefix)
    if out and out[-1].strip():
        out.append("")
    out.append(block)
    if suffix:
        out.append("")
        out.extend(suffix)
    return "\n".join(out).rstrip()


def process_lecture(lecture_file: str, targets: list[Target]) -> None:
    path = LECTURE_ROOT / lecture_file
    markdown = path.read_text(encoding="utf-8").rstrip() + "\n"
    slides = split_slides(markdown)

    for target in targets:
        image_paths = extract_images(target.source_lecture, target.source_slide)
        if not image_paths:
            print(f"warning: no images in {target.source_lecture} slide {target.source_slide}")
            continue
        idx = find_slide_index(slides, target)
        alt_prefix = f"{target.source_lecture} slide {target.source_slide:02d}"
        slides[idx] = inject_photos(slides[idx], image_paths, alt_prefix)
        print(f"updated {lecture_file} slide_index={idx+1} with {len(image_paths)} images from {target.source_lecture}:{target.source_slide}")

    path.write_text("\n---\n".join(slides).rstrip() + "\n", encoding="utf-8")


def main() -> None:
    targets_by_file: dict[str, list[Target]] = {}
    for target in TARGETS:
        targets_by_file.setdefault(target.lecture_file, []).append(target)

    for lecture_file, file_targets in targets_by_file.items():
        process_lecture(lecture_file, file_targets)


if __name__ == "__main__":
    main()
